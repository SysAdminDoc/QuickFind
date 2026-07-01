"""Pluggable text extraction adapters for content search."""

import logging
import os
import importlib.util
import re
import sys
from email import policy
from email.parser import BytesParser
from html import unescape
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol


logger = logging.getLogger('QuickFind.Content')

MAX_EXTRACT_CHARS = 1_000_000
MAX_TEXT_BYTES = 10 * 1024 * 1024

TEXT_EXTENSIONS = {
    'txt', 'md', 'log', 'ini', 'cfg', 'conf', 'json', 'xml', 'yaml', 'yml',
    'csv', 'tsv', 'html', 'htm', 'css', 'scss', 'sass', 'less', 'js',
    'jsx', 'mjs', 'cjs', 'ts', 'tsx', 'vue', 'svelte', 'py', 'pyw',
    'ps1', 'psm1', 'psd1', 'bat', 'cmd', 'sh', 'bash', 'zsh', 'fish',
    'c', 'cc', 'cpp', 'cxx', 'h', 'hpp', 'hh', 'cs', 'fs', 'java',
    'kt', 'kts', 'rs', 'go', 'rb', 'php', 'sql', 'r', 'lua', 'pl',
    'swift', 'scala', 'dart', 'erl', 'ex', 'exs', 'clj', 'cljs',
    'toml', 'env', 'gitignore', 'dockerfile', 'makefile', 'gradle',
}

WINDOWS_SEARCH_EXTENSIONS = {
    'pdf', 'doc', 'docx', 'dot', 'dotx', 'rtf',
    'xls', 'xlsx', 'xlsm', 'xlsb', 'ppt', 'pptx', 'pps', 'ppsx',
    'msg', 'eml', 'odt', 'ods', 'odp', 'one', 'vsd', 'vsdx', 'xps',
}


@dataclass(frozen=True)
class ExtractedContent:
    text: str
    extractor: str


@dataclass(frozen=True)
class AdapterDiagnostic:
    name: str
    extensions: tuple[str, ...]
    available: bool
    detail: str = ""


class ContentAdapter(Protocol):
    name: str
    extensions: set[str]

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        ...


class PlainTextAdapter:
    name = 'text'
    extensions = TEXT_EXTENSIONS

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        if os.path.getsize(path) > MAX_TEXT_BYTES:
            max_chars = min(max_chars, MAX_TEXT_BYTES)
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(max_chars)


class PdfAdapter:
    name = 'pdfplumber'
    extensions = {'pdf'}

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        import pdfplumber

        chunks = []
        total = 0
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ''
                if not text:
                    continue
                remaining = max_chars - total
                if remaining <= 0:
                    break
                chunks.append(text[:remaining])
                total += len(chunks[-1])
        if not chunks:
            return _ocr_pdf_with_tesseract(path, max_chars=max_chars)
        return '\n'.join(chunks)


class DocxAdapter:
    name = 'python-docx'
    extensions = {'docx'}

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        from docx import Document

        doc = Document(path)
        chunks = []
        total = 0
        for paragraph in doc.paragraphs:
            text = paragraph.text
            if not text:
                continue
            remaining = max_chars - total
            if remaining <= 0:
                break
            chunks.append(text[:remaining])
            total += len(chunks[-1])
        return '\n'.join(chunks)


class PptxAdapter:
    name = 'python-pptx'
    extensions = {'pptx'}

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        from pptx import Presentation

        deck = Presentation(path)
        chunks = []
        total = 0
        for slide in deck.slides:
            for shape in slide.shapes:
                text = getattr(shape, 'text', '')
                if not text:
                    continue
                remaining = max_chars - total
                if remaining <= 0:
                    return '\n'.join(chunks)
                chunks.append(text[:remaining])
                total += len(chunks[-1])
        return '\n'.join(chunks)


class EmlAdapter:
    name = 'eml'
    extensions = {'eml'}

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        with open(path, 'rb') as f:
            message = BytesParser(policy=policy.default).parse(f)

        chunks = []
        for header in ("From", "To", "Cc", "Bcc", "Subject", "Date"):
            value = message.get(header)
            if value:
                chunks.append(f"{header}: {value}")

        plain_parts = []
        html_parts = []
        parts = message.walk() if message.is_multipart() else [message]
        for part in parts:
            if part.get_content_maintype() == 'multipart':
                continue
            if part.get_content_disposition() == 'attachment':
                continue
            content_type = part.get_content_type()
            try:
                payload = part.get_content()
            except Exception:
                payload = ""
            if not isinstance(payload, str):
                continue
            if content_type == 'text/plain':
                plain_parts.append(payload)
            elif content_type == 'text/html':
                html_parts.append(_html_to_text(payload))

        chunks.extend(plain_parts or html_parts)
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text[:max_chars]


class WindowsSearchAdapter:
    """Use Windows Search indexed properties backed by installed IFilters."""

    name = 'windows-ifilter'
    extensions = WINDOWS_SEARCH_EXTENSIONS

    @classmethod
    def availability(cls) -> tuple[bool, str]:
        if sys.platform != 'win32':
            return False, "Windows Search COM APIs are only available on Windows"
        missing = []
        for module_name in ("pythoncom", "win32com.client"):
            try:
                spec = importlib.util.find_spec(module_name)
            except (ImportError, ValueError):
                spec = None
            if spec is None:
                missing.append(module_name)
        if missing:
            return False, "Missing Python module: " + ", ".join(missing)
        return True, "Uses Windows Search property handlers/IFilters for indexed files"

    def extract(self, path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
        pythoncom, win32_client = _load_windows_search_com()
        initialized = False
        connection = None
        recordset = None
        try:
            co_initialize = getattr(pythoncom, "CoInitialize", None)
            if callable(co_initialize):
                co_initialize()
                initialized = True
            connection = win32_client.Dispatch("ADODB.Connection")
            connection.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")
            recordset = connection.Execute(_windows_search_sql(path))
            if isinstance(recordset, tuple):
                recordset = recordset[0]
            return _windows_search_record_text(recordset, max_chars)
        finally:
            _close_com_object(recordset)
            _close_com_object(connection)
            co_uninitialize = getattr(pythoncom, "CoUninitialize", None)
            if initialized and callable(co_uninitialize):
                co_uninitialize()


ADAPTERS: tuple[ContentAdapter, ...] = (
    PlainTextAdapter(),
    PdfAdapter(),
    DocxAdapter(),
    PptxAdapter(),
    EmlAdapter(),
    WindowsSearchAdapter(),
)

_ADAPTERS_BY_EXTENSION: dict[str, list[ContentAdapter]] = {}
for adapter in ADAPTERS:
    for ext in adapter.extensions:
        _ADAPTERS_BY_EXTENSION.setdefault(ext, []).append(adapter)

SUPPORTED_CONTENT_EXTENSIONS = set(_ADAPTERS_BY_EXTENSION)


def is_supported_content_path(path: str) -> bool:
    return _extension(path) in SUPPORTED_CONTENT_EXTENSIONS


def adapter_for_path(path: str) -> ContentAdapter | None:
    adapters = adapters_for_path(path)
    return adapters[0] if adapters else None


def adapters_for_path(path: str, available_only: bool = True) -> tuple[ContentAdapter, ...]:
    adapters = _ADAPTERS_BY_EXTENSION.get(_extension(path), [])
    if not available_only:
        return tuple(adapters)
    return tuple(adapter for adapter in adapters if _adapter_available_detail(adapter)[0])


def adapter_diagnostics() -> list[AdapterDiagnostic]:
    diagnostics = []
    for adapter in ADAPTERS:
        available, detail = _adapter_available_detail(adapter)
        diagnostics.append(AdapterDiagnostic(
            name=adapter.name,
            extensions=tuple(sorted(adapter.extensions)),
            available=available,
            detail=detail,
        ))
    diagnostics.append(_ocr_diagnostic())
    return diagnostics


def extract_text(path: str, max_chars: int = MAX_EXTRACT_CHARS) -> ExtractedContent | None:
    adapters = adapters_for_path(path)
    if not adapters:
        return None
    empty_result = None
    for adapter in adapters:
        try:
            text = adapter.extract(path, max_chars=max_chars)
        except (OSError, PermissionError, UnicodeError, ImportError) as exc:
            logger.debug(f"Content extraction failed for {path} with {adapter.name}: {exc}")
            continue
        except Exception as exc:
            logger.debug(f"Content extraction failed for {path} with {adapter.name}: {exc}")
            continue
        result = ExtractedContent(text=text, extractor=adapter.name)
        if text:
            return result
        empty_result = result
    return empty_result


def matched_line_context(content: str, search_text: str,
                         case_sensitive: bool = False,
                         context_lines: int = 3,
                         max_matches: int = 5) -> str:
    if not content or not search_text:
        return content

    needle = search_text if case_sensitive else search_text.lower()
    lines = content.splitlines()

    matches = []
    for idx, line in enumerate(lines):
        segment = line if case_sensitive else line.lower()
        if needle in segment:
            matches.append(idx)
            if len(matches) >= max_matches:
                break

    if not matches:
        return content

    ranges = []
    for idx in matches:
        start = max(0, idx - context_lines)
        end = min(len(lines), idx + context_lines + 1)
        if ranges and start <= ranges[-1][1]:
            ranges[-1] = (ranges[-1][0], max(ranges[-1][1], end))
        else:
            ranges.append((start, end))

    output = []
    for range_index, (start, end) in enumerate(ranges):
        if range_index:
            output.append("...")
        for line_number in range(start, end):
            prefix = ">" if line_number in matches else " "
            output.append(f"{prefix} {line_number + 1}: {lines[line_number]}")
    return '\n'.join(output)


def _extension(path: str) -> str:
    suffix = Path(path).suffix.lower().lstrip('.')
    if suffix:
        return suffix
    return Path(path).name.lower()


def _html_to_text(html: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = unescape(text)
    return re.sub(r"\s+", " ", text).strip()


def _adapter_available_detail(adapter: ContentAdapter) -> tuple[bool, str]:
    availability = getattr(adapter, "availability", None)
    if callable(availability):
        return availability()
    module_names = {
        'text': None,
        'pdfplumber': 'pdfplumber',
        'python-docx': 'docx',
        'python-pptx': 'pptx',
        'eml': None,
    }
    module_name = module_names.get(adapter.name)
    if not module_name:
        return True, ""
    available = importlib.util.find_spec(module_name) is not None
    detail = "" if available else f"Missing Python module: {module_name}"
    return available, detail


def _load_windows_search_com():
    import pythoncom
    import win32com.client

    return pythoncom, win32com.client


def _windows_search_sql(path: str) -> str:
    query_path = os.path.abspath(path) if sys.platform == 'win32' else path
    escaped_path = query_path.replace("'", "''")
    columns = ", ".join([
        "System.Search.Contents",
        "System.Title",
        "System.Subject",
        "System.Author",
        "System.Keywords",
        "System.Comment",
    ])
    return (
        f"SELECT {columns} "
        "FROM SYSTEMINDEX "
        f"WHERE System.ItemPathDisplay = '{escaped_path}'"
    )


def _windows_search_record_text(recordset, max_chars: int) -> str:
    if recordset is None or bool(getattr(recordset, "EOF", True)):
        return ""
    chunks = []
    contents = _record_field_text(recordset, "System.Search.Contents")
    if contents:
        chunks.append(contents)
    metadata_fields = [
        ("Title", "System.Title"),
        ("Subject", "System.Subject"),
        ("Author", "System.Author"),
        ("Keywords", "System.Keywords"),
        ("Comment", "System.Comment"),
    ]
    for label, field_name in metadata_fields:
        value = _record_field_text(recordset, field_name)
        if value:
            chunks.append(f"{label}: {value}")
    return "\n".join(chunks)[:max_chars]


def _record_field_text(recordset, field_name: str) -> str:
    fields = getattr(recordset, "Fields", None)
    if fields is None:
        return ""
    field = None
    try:
        item = getattr(fields, "Item", None)
        if callable(item):
            field = item(field_name)
        elif callable(fields):
            field = fields(field_name)
    except Exception:
        field = None
    if field is None:
        return ""
    value = getattr(field, "Value", None)
    return _plain_field_value(value)


def _plain_field_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, bytes):
        value = value.decode("utf-8", errors="replace")
    elif isinstance(value, (list, tuple, set)):
        value = ", ".join(_plain_field_value(item) for item in value if item is not None)
    else:
        value = str(value)
    return re.sub(r"\s+", " ", value).strip()


def _close_com_object(value) -> None:
    close = getattr(value, "Close", None)
    if callable(close):
        try:
            close()
        except Exception:
            pass


def _ocr_diagnostic() -> AdapterDiagnostic:
    missing = []
    if importlib.util.find_spec("pytesseract") is None:
        missing.append("pytesseract")
    if importlib.util.find_spec("pypdfium2") is None:
        missing.append("pypdfium2")
    if missing:
        return AdapterDiagnostic(
            name="tesseract-ocr",
            extensions=("pdf",),
            available=False,
            detail="Missing optional module: " + ", ".join(missing),
        )
    return AdapterDiagnostic(
        name="tesseract-ocr",
        extensions=("pdf",),
        available=True,
        detail="Requires a working Tesseract executable on PATH.",
    )


def _ocr_pdf_with_tesseract(path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    try:
        import pypdfium2 as pdfium
        import pytesseract
    except ImportError:
        return ""

    chunks = []
    total = 0
    pdf = pdfium.PdfDocument(path)
    try:
        for page in pdf:
            remaining = max_chars - total
            if remaining <= 0:
                break
            try:
                image = page.render(scale=2).to_pil()
                text = pytesseract.image_to_string(image) or ""
            except Exception as exc:
                logger.debug("PDF OCR failed for %s: %s", path, exc)
                continue
            if not text:
                continue
            chunks.append(text[:remaining])
            total += len(chunks[-1])
    finally:
        close = getattr(pdf, "close", None)
        if callable(close):
            close()
    return "\n".join(chunks)
